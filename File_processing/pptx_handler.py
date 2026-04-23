import sys
import os
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from config import DEBUG, GROQ_API_KEY, GROQ_VISION_MODEL
from groq import Groq
from pathlib import Path
import base64
import re

groq_client = Groq(api_key=GROQ_API_KEY)

# ── PLACEHOLDER TOKENS ────────────────────────────────
_PLACEHOLDER_TOKENS = {"‹#›", "<#>", "‹›", "#", "​"}


# ── RECURSIVE SHAPE ITERATOR ──────────────────────────
def iter_shapes(shapes):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


# ── TITLE DETECTION ───────────────────────────────────
def is_title_shape(shape) -> bool:
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER
        if shape.is_placeholder:
            ph = shape.placeholder_format.type
            return ph in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
    except Exception:
        pass
    return shape.name.lower().startswith("title")


# ── TEXT CLEANING ─────────────────────────────────────
def clean_content(text: str) -> str:
    """Strip placeholder tokens and collapse extra blank lines."""
    for token in _PLACEHOLDER_TOKENS:
        text = text.replace(token, "")
    return re.sub(r'\n{3,}', '\n\n', text).strip()


# ── SLIDE TEXT EXTRACTION ─────────────────────────────
def extract_slide_text(sorted_shapes) -> str:
    """Extract all text from a slide's shapes, title gets # prefix."""
    slide_text = ""
    for shape in sorted_shapes:
        if shape.has_text_frame:
            prefix = "# " if is_title_shape(shape) else ""
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    slide_text += prefix + text + "\n\n"
                    prefix = ""
    return clean_content(slide_text)


# ── TABLE EXTRACTION ──────────────────────────────────
def extract_slide_tables(sorted_shapes) -> list[str]:
    """Extract all tables from a slide as markdown strings."""
    tables = []
    for shape in sorted_shapes:
        if shape.has_table:
            try:
                rows = []
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(cells)

                if rows:
                    headers  = rows[0]
                    data     = rows[1:]
                    md_table = "| " + " | ".join(headers) + " |\n"
                    md_table += "| " + " | ".join(["---"] * len(headers)) + " |\n"
                    for row in data:
                        while len(row) < len(headers):
                            row.append("")
                        md_table += "| " + " | ".join(row) + " |\n"
                    tables.append(md_table)

            except Exception as e:
                if DEBUG:
                    print(f"[PPTX] Table error: {e}")
    return tables


# ── IMAGE EXTRACTION ──────────────────────────────────
def extract_slide_images(slide) -> list[bytes]:
    """Extract raw image blobs from a slide."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    _IMAGE_TYPES = {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE, MSO_SHAPE_TYPE.MEDIA}
    images = []
    for shape in iter_shapes(slide.shapes):
        try:
            if shape.shape_type in _IMAGE_TYPES or hasattr(shape, "image"):
                images.append(shape.image.blob)
        except Exception as e:
            if DEBUG:
                print(f"[PPTX] Image extraction error: {e}")
    return images


def describe_image(img_bytes: bytes, current_section: str) -> str:
    """Send image to Groq Vision and return a text description."""
    img_base64   = base64.b64encode(img_bytes).decode("utf-8")
    context_hint = f"This image is from a presentation slide in the section '{current_section}'. " if current_section else ""

    response = groq_client.chat.completions.create(
        model=GROQ_VISION_MODEL,
        messages=[{
            "role": "user",
            "content": [
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:image/png;base64,{img_base64}"}
                },
                {
                    "type": "text",
                    "text": f"{context_hint}Describe this image briefly."
                }
            ]
        }],
        max_tokens=200,
    )
    return response.choices[0].message.content


# ── SECTION DETECTION HELPERS ─────────────────────────
def total_paragraph_count(sorted_shapes) -> int:
    """Count all non-empty paragraphs across every shape on the slide."""
    total = 0
    for shape in sorted_shapes:
        if shape.has_text_frame:
            total += sum(
                1 for p in shape.text_frame.paragraphs
                if p.text.strip()
            )
    return total


def has_bullet_points(sorted_shapes) -> bool:
    """True if any shape contains bullet-style paragraphs."""
    for shape in sorted_shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            if para.level > 0:
                return True
            text = para.text.strip()
            if text and text[0] in ("•", "→", "✓", "-", "▶", "◆", "★"):
                return True
    return False


def total_word_count(sorted_shapes) -> int:
    """Count total words across all shapes on the slide."""
    total_text = " ".join(
        s.text_frame.text.strip()
        for s in sorted_shapes
        if s.has_text_frame
    )
    return len(total_text.split())


def is_section_title_slide(slide, sorted_shapes) -> bool:
    """Returns True if this slide is a section divider with no real content."""

    # Hard disqualifiers
    has_table = any(s.has_table for s in iter_shapes(slide.shapes))
    if has_table:
        return False

    from pptx.enum.shapes import MSO_SHAPE_TYPE
    _IMAGE_TYPES = {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE, MSO_SHAPE_TYPE.MEDIA}
    has_image = any(
        s.shape_type in _IMAGE_TYPES or hasattr(s, "image")
        for s in iter_shapes(slide.shapes)
    )
    if has_image:
        return False

    # Content volume checks
    if total_paragraph_count(sorted_shapes) > 4:
        return False
    if has_bullet_points(sorted_shapes):
        return False
    if total_word_count(sorted_shapes) > 30:
        return False

    # Must have a title shape
    if not any(is_title_shape(s) and s.has_text_frame for s in sorted_shapes):
        return False

    # Title itself must be short
    for shape in sorted_shapes:
        if is_title_shape(shape) and shape.has_text_frame:
            title_text  = shape.text_frame.text.strip()
            title_lines = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
            if len(title_text) > 80 or len(title_lines) > 3:
                return False

    return True


def get_section_name(slide, sorted_shapes, slide_num: int) -> str:
    title_texts = []   # ← collect ALL title shapes
    other_texts = []

    for shape in sorted_shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text or text in _PLACEHOLDER_TOKENS:
            continue
        if is_title_shape(shape):
            title_texts.append(text)  # ← append, don't overwrite
        else:
            other_texts.append(text)

    # Filter noise (page numbers, single chars) from other_texts
    other_texts = [
        t for t in other_texts
        if not t.strip().isdigit() and len(t.strip()) > 2
    ]

    # Separate numeric labels ("01", "02") from real titles
    numeric_titles = [t for t in title_texts if t.strip().isdigit() or len(t.strip()) <= 3]
    real_titles    = [t for t in title_texts if t not in numeric_titles]

    # Best real title = longest non-numeric title shape
    if real_titles:
        best_title = max(real_titles, key=len)
        if numeric_titles:
            return f"{numeric_titles[0]} - {best_title}"  # "02 - Présentation du dataset"
        return best_title                                  # "Introduction"

    # No real title found → fall back to other_texts
    if numeric_titles and other_texts:
        return f"{numeric_titles[0]} - {max(other_texts, key=len)}"

    # Last resort
    return title_texts[0] if title_texts else f"Section {slide_num}"

# ── SLIDE PARSER ──────────────────────────────────────
def parse_pptx(file_path: str) -> list[dict]:
    """
    Parse a PPTX file and return a structured dict of slides.

    Returns:
        {
            "filename": str,
            "total_slides": int,
            "slides": [
                {
                    "slide_num":    int,
                    "is_section":   bool,
                    "section_name": str | None,   # only if is_section=True
                    "current_section": str | None, # inherited section label
                    "text":         str,
                    "tables":       [str, ...],    # markdown tables
                    "images":       [str, ...],    # Groq descriptions
                }
            ]
        }
    """
    from pptx import Presentation

    if DEBUG:
        print(f"\n[PPTX] Parsing: {file_path}")

    prs            = Presentation(file_path)
    filename       = Path(file_path).name
    current_section = None
    slides_data    = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        if DEBUG:
            print(f"[PPTX] Slide {slide_num}...")

        sorted_shapes = sorted(
            iter_shapes(slide.shapes),
            key=lambda s: (
                s.top  if s.top  is not None else 0,
                s.left if s.left is not None else 0,
            )
        )

        is_section = is_section_title_slide(slide, sorted_shapes)
        section_name = None

        if is_section:
            section_name    = get_section_name(slide, sorted_shapes, slide_num)
            current_section = section_name
            if DEBUG:
                print(f"[PPTX] Section detected: {current_section}")
        elif DEBUG:
            print(f"[PPTX] Content slide (section={current_section})")

        # Extract content
        text   = extract_slide_text(sorted_shapes)
        tables = extract_slide_tables(sorted_shapes)
        images = []

        img_blobs = extract_slide_images(slide)
        if DEBUG:
            print(f"[PPTX] Found {len(img_blobs)} images")

        for img_bytes in img_blobs:
            try:
                description = describe_image(img_bytes, current_section)
                images.append(description)
            except Exception as e:
                if DEBUG:
                    print(f"[PPTX] Image error: {e}")

        slides_data.append({
            "slide_num":       slide_num,
            "is_section":      is_section,
            "section_name":    section_name,
            "current_section": current_section,
            "text":            text,
            "tables":          tables,
            "images":          images,
        })

    return {
        "filename":     filename,
        "total_slides": len(prs.slides),
        "slides":       slides_data,
    }
# ── TEST ──────────────────────────────────────────────
if __name__ == "__main__":
    import json

    parsed = parse_pptx("2.pptx")

    output_path = Path("2.pptx").stem + "_handler_output.json"
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(parsed, f, ensure_ascii=False, indent=2)

    print(f"Saved to {output_path}")
    print(f"Total slides : {parsed['total_slides']}")
    print(f"Sections     : {[s['section_name'] for s in parsed['slides'] if s['is_section']]}")
    print(f"Content slides: {[s['slide_num'] for s in parsed['slides'] if not s['is_section']]}")