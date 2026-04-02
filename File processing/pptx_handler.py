import sys
import os
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from config import DEBUG, GROQ_API_KEY, GROQ_VISION_MODEL
from groq import Groq
from pathlib import Path
import base64
import json

groq_client = Groq(api_key=GROQ_API_KEY)


def extract_pptx(file_path: str) -> dict:

    if DEBUG:
        print(f"\n[PPTX] Processing: {file_path}")

    from pptx import Presentation

    prs      = Presentation(file_path)
    filename = Path(file_path).name
    chunks   = []
    img_counter = 0
    tbl_counter = 0

    for slide_num, slide in enumerate(prs.slides, start=1):

        if DEBUG:
            print(f"[PPTX] Slide {slide_num}...")

        # ── SORT SHAPES BY POSITION ───────────────────
        sorted_shapes = sorted(
            slide.shapes,
            key=lambda s: (
                s.top  if s.top  is not None else 0,
                s.left if s.left is not None else 0,
            )
        )

        # ── DETECT SLIDE TITLE / SECTION ─────────────
        section_title = None
        for shape in sorted_shapes:
            if shape.has_text_frame and shape.name.lower().startswith("title"):
                section_title = shape.text_frame.text.strip()
                break
        if not section_title:
            section_title = f"Slide {slide_num}"

        slide_text = ""

        for shape in sorted_shapes:

            # ── TEXT ──────────────────────────────────
            if shape.has_text_frame:
                # skip title shape — already captured as section
                if shape.name.lower().startswith("title"):
                    continue
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_text += text + "\n\n"

            # ── TABLE ─────────────────────────────────
            if shape.has_table:
                try:
                    table = shape.table
                    rows  = []

                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(cells)

                    if rows:
                        headers  = rows[0]
                        data     = rows[1:]
                        md_table = "| " + " | ".join(headers) + " |\n"
                        md_table += "| " + " | ".join(
                            ["---"] * len(headers)
                        ) + " |\n"
                        for row in data:
                            while len(row) < len(headers):
                                row.append("")
                            md_table += "| " + " | ".join(row) + " |\n"

                        tbl_counter += 1
                        chunks.append({
                            "type"    : "table",
                            "content" : md_table,
                            "metadata": {
                                "source"     : filename,
                                "slide"      : slide_num,
                                "section"    : section_title,
                                "table_index": tbl_counter,
                            }
                        })

                        if DEBUG:
                            print(f"[PPTX] Table {tbl_counter} "
                                  f"(slide {slide_num}) ")

                except Exception as e:
                    if DEBUG:
                        print(f"[PPTX] Table error: {e}")

        # ── FLUSH SLIDE TEXT AS ONE CHUNK ─────────────
        if slide_text.strip():
            chunks.append({
                "type"    : "text",
                "content" : slide_text.strip(),
                "metadata": {
                    "source" : filename,
                    "slide"  : slide_num,
                    "section": section_title,
                }
            })

        # ── IMAGES ────────────────────────────────────
        for img_bytes in _extract_images_from_slide(slide):
            try:
                img_base64 = base64.b64encode(img_bytes).decode("utf-8")

                response = groq_client.chat.completions.create(
                    model=GROQ_VISION_MODEL,
                    messages=[{
                        "role": "user",
                        "content": [
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/png;base64,{img_base64}"
                                }
                            },
                            {
                                "type": "text",
                                "text": (
                                    "Describe this image briefly. "
                                    "Focus only on meaningful visual content."
                                )
                            }
                        ]
                    }],
                    max_tokens=200
                )

                description = response.choices[0].message.content
                img_counter += 1

                chunks.append({
                    "type"    : "image",
                    "content" : description,
                    "metadata": {
                        "source"     : filename,
                        "slide"      : slide_num,
                        "section"    : section_title,
                        "image_index": img_counter,
                    }
                })

                if DEBUG:
                    print(f"[PPTX] Image {img_counter} "
                          f"(slide {slide_num}) ")

            except Exception as e:
                if DEBUG:
                    print(f"[PPTX] Image error: {e}")

    result = {
        "filename": filename,
        "chunks"  : chunks,
        "metadata": {
            "source"     : filename,
            "slide_count": len(prs.slides),
            "text_count" : len([c for c in chunks if c["type"] == "text"]),
            "table_count": len([c for c in chunks if c["type"] == "table"]),
            "image_count": len([c for c in chunks if c["type"] == "image"]),
        }
    }

    # ── SAVE AS JSON ──────────────────────────────────
    output_path = Path(file_path).stem + "_parsed.json"
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(result, f, ensure_ascii=False, indent=2)

    if DEBUG:
        print(f"\n[PPTX] Slides : {len(prs.slides)}")
        print(f"[PPTX] Text   : {result['metadata']['text_count']}")
        print(f"[PPTX] Tables : {result['metadata']['table_count']}")
        print(f"[PPTX] Images : {result['metadata']['image_count']}")
        print(f"[PPTX] Saved  : {output_path}")

    return result


# ── IMAGE EXTRACTION ──────────────────────────────────
def _extract_images_from_slide(slide) -> list:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    images = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                if shape.width < 500000 or shape.height < 500000:
                    continue
                images.append(shape.image.blob)
            except Exception as e:
                if DEBUG:
                    print(f"[PPTX] Image extraction error: {e}")
    return images


# ── TEST ──────────────────────────────────────────────
if __name__ == "__main__":
    result = extract_pptx("../test9.pptx")

    print("\n===== CHUNKS PREVIEW =====")
    for i, chunk in enumerate(result["chunks"][:5]):
        print(f"\nChunk {i+1}:")
        print(f"  Type   : {chunk['type']}")
        print(f"  Slide  : {chunk['metadata']['slide']}")
        print(f"  Section: {chunk['metadata']['section']}")
        print(f"  Content: {chunk['content'][:100]}...")

    print(f"\n Total chunks : {len(result['chunks'])}")
    print(f" Saved to     : {Path('../test3').stem}_parsed.json")